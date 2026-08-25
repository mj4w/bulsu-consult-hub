from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


DOWNLOADS = Path.home() / "Downloads"
SOURCE = DOWNLOADS / "My presentation - softengr.pptx"
BACKUP = DOWNLOADS / "My presentation - softengr - original backup.pptx"
OUTPUT = DOWNLOADS / "My presentation - softengr - edited.pptx"
SCRIPT_OUTPUT = DOWNLOADS / "AI-Assisted Software Development and Testing - speaking script.txt"

W = Inches(13.333333)
H = Inches(7.5)

NAVY = RGBColor(20, 38, 62)
TEAL = RGBColor(23, 132, 120)
BLUE = RGBColor(42, 111, 197)
GOLD = RGBColor(227, 167, 52)
RED = RGBColor(196, 73, 65)
INK = RGBColor(33, 43, 54)
MUTED = RGBColor(98, 110, 125)
PALE = RGBColor(245, 248, 250)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text, size=22, bold=False, color=INK, align=None):
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in tf.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=22, bold=False, color=INK, align=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    set_text(shape, text, size=size, bold=bold, color=color, align=align)
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_header(slide, section, title):
    add_round_rect(slide, 0, 0, Inches(3.1), H, NAVY, radius=False)
    add_round_rect(slide, Inches(0.7), Inches(0.75), Inches(1.0), Inches(1.0), TEAL)
    add_textbox(slide, Inches(0.72), Inches(0.93), Inches(0.96), Inches(0.35), "AI", 22, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(2.05), Inches(2.0), Inches(0.35), section, 15, True, RGBColor(130, 216, 205))
    add_textbox(slide, Inches(0.7), Inches(2.48), Inches(2.1), Inches(1.7), title, 30, True, WHITE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.23), Inches(1.25), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_title_slide(slide, section, title, subtitle):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_round_rect(slide, Inches(0.78), Inches(0.75), Inches(1.0), Inches(0.55), TEAL)
    add_textbox(slide, Inches(0.92), Inches(0.88), Inches(0.7), Inches(0.22), section, 14, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.78), Inches(1.65), Inches(8.9), Inches(1.25), title, 38, True, NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(3.05), Inches(2.55), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    add_textbox(slide, Inches(0.78), Inches(3.38), Inches(8.6), Inches(1.0), subtitle, 22, False, MUTED)
    add_round_rect(slide, Inches(9.2), Inches(0.85), Inches(2.95), Inches(5.6), WHITE, RGBColor(220, 226, 232))
    for i, label in enumerate(["Generative AI", "Agentic AI", "Vibe coding"]):
        y = Inches(1.35 + i * 1.45)
        add_round_rect(slide, Inches(9.58), y, Inches(2.18), Inches(0.68), [BLUE, TEAL, GOLD][i])
        add_textbox(slide, Inches(9.72), y + Inches(0.18), Inches(1.9), Inches(0.22), label, 14, True, WHITE, PP_ALIGN.CENTER)


def add_bullet_row(slide, x, y, text, color=TEAL):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_textbox(slide, x + Inches(0.38), y, Inches(5.95), Inches(0.58), text, 20, False, INK)


def add_three_column_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_textbox(slide, Inches(0.78), Inches(0.55), Inches(10.6), Inches(0.45), "2.9 - AI CONCEPTS", 15, True, TEAL)
    add_textbox(slide, Inches(0.78), Inches(0.95), Inches(10.8), Inches(0.7), "Generative AI, Agentic AI, and AI Slop", 31, True, NAVY)
    items = [
        ("Generative AI", BLUE, "Creates software artifacts: code, documentation, diagrams, explanations, and test cases from prompts or project context."),
        ("Agentic AI", TEAL, "Plans multi-step work: reads files, uses tools, edits code, runs checks, and iterates toward a development goal."),
        ("AI Slop", RED, "Low-quality AI output: confident-looking but shallow, incorrect, generic, insecure, or untested work that creates cleanup cost."),
    ]
    for i, (heading, color, body) in enumerate(items):
        x = Inches(0.78 + i * 4.12)
        add_round_rect(slide, x, Inches(2.0), Inches(3.55), Inches(3.85), WHITE, RGBColor(219, 225, 232))
        add_round_rect(slide, x + Inches(0.28), Inches(2.32), Inches(0.62), Inches(0.62), color)
        add_textbox(slide, x + Inches(1.05), Inches(2.26), Inches(2.12), Inches(0.5), heading, 20, True, NAVY)
        add_textbox(slide, x + Inches(0.34), Inches(3.25), Inches(2.88), Inches(1.72), body, 18, False, INK)
        add_textbox(slide, x + Inches(0.34), Inches(5.15), Inches(2.85), Inches(0.35), "Key idea: verify before trusting.", 14, True, MUTED)
    add_textbox(slide, Inches(11.78), Inches(6.74), Inches(0.55), Inches(0.25), "21", 12, False, MUTED, PP_ALIGN.RIGHT)


def add_vibe_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "2.10 - CODING STYLE", "Vibe Coding vs. AI-Assisted Development")
    x = Inches(3.7)
    add_textbox(slide, x, Inches(0.65), Inches(7.7), Inches(0.6), "They both use AI, but the level of control is different.", 25, True, NAVY)
    add_round_rect(slide, x, Inches(1.55), Inches(3.6), Inches(4.55), PALE, RGBColor(219, 225, 232))
    add_round_rect(slide, Inches(7.65), Inches(1.55), Inches(3.6), Inches(4.55), PALE, RGBColor(219, 225, 232))
    add_textbox(slide, x + Inches(0.32), Inches(1.9), Inches(2.9), Inches(0.42), "Vibe Coding", 23, True, GOLD)
    add_textbox(slide, x + Inches(0.32), Inches(2.55), Inches(2.75), Inches(2.55), "A developer describes the desired result and lets AI generate large parts quickly. Useful for prototypes and experiments, but risky when the developer cannot explain or maintain the code.", 18, False, INK)
    add_textbox(slide, Inches(7.97), Inches(1.9), Inches(2.9), Inches(0.42), "AI-Assisted Development", 23, True, TEAL)
    add_textbox(slide, Inches(7.97), Inches(2.55), Inches(2.75), Inches(2.55), "The developer stays in control: defines requirements, reviews each change, runs tests, checks security, and accepts only what fits the system.", 18, False, INK)
    add_textbox(slide, x, Inches(6.35), Inches(7.8), Inches(0.35), "Best practice: use AI for speed, but keep engineering judgment, testing, and accountability human-led.", 17, True, NAVY)
    add_textbox(slide, Inches(11.78), Inches(6.74), Inches(0.55), Inches(0.25), "22", 12, False, MUTED, PP_ALIGN.RIGHT)


def add_risk_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_textbox(slide, Inches(0.78), Inches(0.55), Inches(10.0), Inches(0.45), "2.11 - PRODUCTIVITY AND RISKS", 15, True, TEAL)
    add_textbox(slide, Inches(0.78), Inches(0.95), Inches(10.7), Inches(0.7), "Productivity Benefits Need Guardrails", 32, True, NAVY)
    add_round_rect(slide, Inches(0.8), Inches(1.95), Inches(5.55), Inches(4.45), WHITE, RGBColor(219, 225, 232))
    add_round_rect(slide, Inches(6.98), Inches(1.95), Inches(5.55), Inches(4.45), WHITE, RGBColor(219, 225, 232))
    add_textbox(slide, Inches(1.18), Inches(2.25), Inches(4.5), Inches(0.4), "What AI improves", 23, True, TEAL)
    add_bullet_row(slide, Inches(1.22), Inches(3.02), "Less boilerplate and faster first drafts", TEAL)
    add_bullet_row(slide, Inches(1.22), Inches(3.78), "Quicker debugging, review, and documentation", TEAL)
    add_bullet_row(slide, Inches(1.22), Inches(4.54), "More test ideas and edge cases", TEAL)
    add_bullet_row(slide, Inches(1.22), Inches(5.30), "More time for design decisions and problem-solving", TEAL)
    add_textbox(slide, Inches(7.36), Inches(2.25), Inches(4.5), Inches(0.4), "What teams must control", 23, True, RED)
    add_bullet_row(slide, Inches(7.40), Inches(3.02), "Incorrect code that looks correct", RED)
    add_bullet_row(slide, Inches(7.40), Inches(3.78), "Security, privacy, and licensing issues", RED)
    add_bullet_row(slide, Inches(7.40), Inches(4.54), "AI slop entering the codebase", RED)
    add_bullet_row(slide, Inches(7.40), Inches(5.30), "Skill loss if developers stop reading the code", RED)
    add_textbox(slide, Inches(11.78), Inches(6.74), Inches(0.55), Inches(0.25), "23", 12, False, MUTED, PP_ALIGN.RIGHT)


def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_slide(
        slide,
        "2.12",
        "Main Takeaway",
        "AI is not replacing software engineering. It is changing the workflow: developers spend less time starting from a blank page and more time reviewing, testing, integrating, and making judgment calls.",
    )
    add_textbox(slide, Inches(0.78), Inches(5.25), Inches(7.9), Inches(0.85), "Good use of AI = clear instructions + human review + automated tests + responsibility for the final output.", 24, True, NAVY)
    add_textbox(slide, Inches(11.78), Inches(6.74), Inches(0.55), Inches(0.25), "24", 12, False, MUTED, PP_ALIGN.RIGHT)


def replace_existing_text(prs):
    replacements = {
        "AI-Assisted Programming": "AI-Assisted Coding",
        "AI coding tools like Copilot, Claude Code, and Cursor work right beside developers": "AI coding tools work beside developers as coding partners, not final decision-makers",
        "They understand context — the code, comments, and names used": "They use project context: files, comments, function names, and coding patterns",
        "They work with many programming languages and frameworks": "They support many languages, frameworks, and development environments",
        "Pair-programming loop: \nsuggest → review → accept or edit": "Practical loop:\nprompt → review → test → accept or revise",
        "Writes full functions or modules from a plain-language request": "Creates functions, classes, APIs, or starter files from a clear prompt",
        "Can build entire files, classes, or starter code": "Best for first drafts, scaffolding, examples, and repeated coding patterns",
        "Guesses and fills in the next lines as you type": "Predicts the next lines while the developer is typing",
        "Cuts down repetitive typing in real time": "Reduces repetitive typing and speeds up small implementation steps",
        "Explains error messages in plain, simple language": "Explains errors, stack traces, and logs in plain language",
        "Suggests likely causes and possible fixes": "Suggests likely causes, fixes, and checks to confirm the issue",
        "Can scan logs to find exactly where things went wrong": "Can scan logs and point to the part of the workflow that failed",
        "Speeds up troubleshooting, especially in code you do not know well": "Speeds up troubleshooting, especially in unfamiliar codebases",
        "Suggests cleaner, more efficient ways to write existing code": "Suggests cleaner, simpler, and more maintainable versions of existing code",
        "Spots messy patterns, repeated code, and unused code": "Spots duplication, unused code, confusing names, and messy patterns",
        "Recommends ways to make code faster and easier to read": "Recommends changes for readability, performance, and maintainability",
        "Helps update and simplify old codebases": "Helps modernize older code while tests protect existing behavior",
        "Flags possible bugs, rule violations, and security issues": "Flags possible bugs, style violations, security issues, and missing checks",
        "Automatically writes review comments and suggestions": "Drafts review comments, explanations, and possible code changes",
        "Helps, but does not replace human reviewers": "Helps reviewers, but does not replace human accountability",
        "Speeds up code reviews inside Continuous Integration (CI) or Continuous Delivery (CD) pipelines": "Speeds up pull requests and CI/CD review workflows",
        "Writes unit and integration tests from the code or specs": "Generates unit, integration, and edge-case tests from code or requirements",
        "Covers more cases, including the ones people might miss": "Suggests boundary cases, invalid inputs, and failure scenarios",
        "Helps build, run, and maintain complete test suites": "Helps build, run, update, and explain automated test suites",
        "Works with Continuous Integration (CI) pipelines to keep testing continuous": "Connects with CI pipelines so tests run continuously on every change",
        "Cuts time spent on repetitive, low-value coding": "Cuts repetitive coding and setup work",
        "Makes debugging and code review faster": "Shortens debugging, review, and documentation cycles",
        "It lets team try out ideas and features faster": "Lets teams prototype ideas and features faster",
        "Frees up developers for bigger-picture problem-solving": "Frees developers for design, architecture, and problem-solving",
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            update_shape_text(shape, replacements)


def update_shape_text(shape, replacements):
    if hasattr(shape, "text") and shape.text:
        for old, new in replacements.items():
            if old in shape.text:
                shape.text = shape.text.replace(old, new)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Aptos"
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            update_shape_text(child, replacements)


def write_script():
    script = """Speaking Script: AI-Assisted Software Development and Testing

Slide 1 - AI-Assisted Software Development and Testing
This part is about how artificial intelligence supports the software development process. AI is not only used to write code. It can assist with coding, code generation and completion, debugging, code review, test-case generation, automated testing, and productivity. The main idea is that AI can make development faster, but developers still need to review, test, and take responsibility for the final output.

Slide 2 - AI-Assisted Coding
AI-assisted coding means the developer still leads the work, but an AI tool acts like a coding partner. Tools such as Copilot, Cursor, and Claude Code can use context from files, comments, function names, and coding patterns. The best workflow is not to accept everything immediately. The practical loop is prompt, review, test, then accept or revise.

Slide 3 - Code Generation and Code Completion
Code generation and code completion are related, but they are not the same. Code generation creates larger outputs, such as a function, class, API route, or starter file, based on a prompt. Code completion is more immediate. It predicts the next lines while the developer is typing. Both reduce repetitive work, but both still need review because generated code can be incomplete or incorrect.

Slide 4 - Debugging and Error Explanation
AI can help developers understand error messages, stack traces, and logs. It can suggest likely causes and possible fixes, which is useful when working with unfamiliar code. However, an AI explanation should be treated as a hypothesis. The developer still needs to reproduce the issue, apply the fix carefully, and confirm that the bug is actually solved.

Slide 5 - Refactoring and Code Optimization
AI can also help improve existing code. It can suggest cleaner structure, remove duplication, rename confusing parts, and recommend more maintainable solutions. This is important because software is not only written once; it is maintained over time. The risk is that a refactor may accidentally change behavior, so automated tests are needed before and after the change.

Slide 6 - AI-Assisted Code Review
AI-assisted code review can flag possible bugs, style violations, missing checks, and security issues. It can also draft review comments or explain why a change may be risky. This speeds up pull requests and CI/CD workflows. But it should support human reviewers, not replace them, because human reviewers understand the project goals, system design, and user impact.

Slide 7 - Test-Case Generation and Automated Testing
AI can generate unit tests, integration tests, and edge cases from source code or requirements. It can suggest invalid inputs, boundary conditions, and failure scenarios that people might miss. It can also help update automated test suites and explain test failures. This supports quality assurance because tests can run continuously in a CI pipeline.

Slide 8 - Productivity Benefits
The main productivity benefit is speed. AI reduces boilerplate, shortens debugging and code review cycles, helps teams prototype faster, and gives developers more time for design and problem-solving. But productivity should not mean blindly accepting output. It should mean faster progress with proper review and testing.

Slide 9 - Example Workflow
A practical workflow starts when the developer describes a needed function or feature. AI generates a first draft. Then AI can propose tests, the developer runs those tests, and review tools flag possible issues. The final step is always a developer decision: accept, edit, or reject the suggestion.

Slide 10 - Generative AI, Agentic AI, and AI Slop
Generative AI creates software artifacts such as code, documentation, diagrams, explanations, and test cases. Agentic AI goes further because it can plan multiple steps, use tools, edit files, run checks, and iterate toward a goal. AI slop is the negative side. It means low-quality AI output that looks confident but is shallow, incorrect, generic, insecure, or untested. This is why verification matters.

Slide 11 - Vibe Coding vs. AI-Assisted Development
Vibe coding is when a developer describes the desired result and lets AI generate large parts of the application quickly. It is useful for prototypes, demos, and exploring ideas. The risk is that the developer may not fully understand the generated code. AI-assisted development is more disciplined. The developer still uses AI, but remains responsible for requirements, architecture, review, tests, security, and maintainability.

Slide 12 - Productivity Benefits and Possible Risks
AI can make teams faster by reducing repetitive work, improving documentation, and suggesting more test cases. But teams must control the risks. Generated code may look correct but fail in edge cases. It may introduce security, privacy, or licensing issues. It may also create AI slop if no one reviews it carefully. The guardrails are code review, automated tests, coding standards, and clear human responsibility.

Slide 13 - Main Takeaway
The main takeaway is that AI is changing software engineering, not replacing it. Developers spend less time starting from a blank page and more time directing, reviewing, testing, and integrating AI output. Good AI use means clear instructions, human review, automated tests, and responsibility for the final product.

Short ending line:
AI is powerful, but software quality still depends on developer judgment. The best results happen when AI handles speed and repetition while humans handle correctness, responsibility, and long-term maintainability.
"""
    SCRIPT_OUTPUT.write_text(script, encoding="utf-8")


def main():
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    if not BACKUP.exists():
        copy2(SOURCE, BACKUP)
    prs = Presentation(SOURCE)
    replace_existing_text(prs)
    add_three_column_slide(prs)
    add_vibe_slide(prs)
    add_risk_slide(prs)
    add_closing_slide(prs)
    prs.save(OUTPUT)
    write_script()
    print(f"Saved {OUTPUT}")
    print(f"Saved {SCRIPT_OUTPUT}")
    print(f"Backup {BACKUP}")


if __name__ == "__main__":
    main()
